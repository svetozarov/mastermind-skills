"""
Universal Gamma font installer v2.
Extracts fonts from PPTX, downloads FULL versions, verifies Cyrillic, installs.
If a font has no Cyrillic — replaces it in PPTX XML with closest standard match.

Usage:
    python install-gamma-fonts.py <path-to-pptx>
"""
import sys
import os
import shutil
import zipfile
import platform
import json
from pathlib import Path
from collections import Counter

sys.stdout.reconfigure(encoding='utf-8')


# --- Standard Windows fonts (no install needed) ---
STANDARD_FONTS = {
    'Calibri', 'Calibri Light', 'Arial', 'Times New Roman',
    'Segoe UI', 'Segoe UI Light', 'Segoe UI Semibold',
    'Tahoma', 'Verdana', 'Cambria', 'Consolas', 'Courier New',
    'Georgia', 'Trebuchet MS', 'Impact', 'Comic Sans MS',
    'Palatino Linotype', 'Book Antiqua', 'Garamond',
    'Century Gothic', 'Franklin Gothic Medium',
    'Candara', 'Corbel', 'Constantia', 'Bahnschrift',
    'Aptos', 'Aptos Display',
    'Malgun Gothic', 'Microsoft YaHei', 'Yu Gothic',
}

# --- Known fonts with full versions on GitHub ---
GITHUB_SOURCES = {
    'Inter': 'https://github.com/rsms/inter/releases/download/v4.1/Inter-4.1.zip',
    'Outfit': 'https://github.com/Outfitio/Outfit-Fonts/archive/refs/heads/main.zip',
    'Montserrat': 'https://github.com/JulietaUla/Montserrat/archive/refs/heads/master.zip',
    'Geist': 'https://github.com/vercel/geist-font/releases/download/1.8.0/geist-font-1.8.0.zip',
    'Manrope': 'https://github.com/sharanda/manrope/releases/download/v4.505/Manrope-v4.505.zip',
    'Space Grotesk': 'https://github.com/nicholasgoss/spaceGrotesk/archive/refs/heads/main.zip',
}

# --- Cyrillic-safe replacements for fonts without Cyrillic ---
NO_CYRILLIC_REPLACEMENTS = {
    'Heebo': 'Calibri',
    'Heebo Light': 'Calibri Light',
    'Heebo Bold': 'Calibri Bold',
}

# --- Cyrillic test codepoints (А-Я, а-я, ё, Ё) ---
CYRILLIC_TEST = [
    0x0410, 0x0411, 0x0412, 0x0413, 0x0414, 0x0415,
    0x0416, 0x0417, 0x0418, 0x0419, 0x041A, 0x041B,
    0x041C, 0x041D, 0x041E, 0x041F, 0x0420,
    0x0430, 0x0431, 0x0432, 0x0433, 0x0451,
]


def extract_fonts_from_pptx(pptx_path):
    """Extract all font names used in a PPTX file."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    fonts = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
    return fonts


def normalize_to_base(font_name):
    """'Montserrat Bold' -> 'Montserrat'"""
    suffixes = [' ExtraBold', ' SemiBold', ' ExtraLight', ' UltraLight',
                ' Bold', ' Light', ' Medium', ' Thin', ' Black',
                ' Italic', ' Regular', ' Display', ' Condensed']
    base = font_name
    for s in suffixes:
        if base.endswith(s):
            base = base[:-len(s)]
    return base.strip() or font_name


def is_standard(font_name):
    """Check if font (or its base) is a standard Windows font."""
    if font_name in STANDARD_FONTS:
        return True
    base = normalize_to_base(font_name)
    return base in STANDARD_FONTS or any(font_name.startswith(std) for std in STANDARD_FONTS)


def download_file(url, dest_path, timeout=60):
    """Download a file from URL."""
    import urllib.request
    req = urllib.request.Request(url, headers={'User-Agent': 'Masterminder/2.0'})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        with open(dest_path, 'wb') as f:
            f.write(resp.read())


def download_from_github(font_name, dest_dir):
    """Download full font from GitHub. Returns list of TTF paths."""
    base = normalize_to_base(font_name)
    if base not in GITHUB_SOURCES:
        return []

    url = GITHUB_SOURCES[base]
    zip_path = dest_dir / f"{base.lower()}-github.zip"

    try:
        print(f"    Downloading from GitHub...")
        download_file(url, zip_path)
    except Exception as e:
        print(f"    [!] GitHub download failed: {e}")
        return []

    ttf_files = []
    try:
        with zipfile.ZipFile(zip_path, 'r') as z:
            for name in z.namelist():
                if name.endswith('.ttf') and 'Variable' not in name and '__MACOSX' not in name:
                    data = z.read(name)
                    out_name = Path(name).name
                    out_path = dest_dir / out_name
                    with open(out_path, 'wb') as f:
                        f.write(data)
                    # Only keep files > 50KB (skip tiny/broken)
                    if out_path.stat().st_size > 50000:
                        ttf_files.append(out_path)
                    else:
                        out_path.unlink()
    except Exception as e:
        print(f"    [!] Unzip error: {e}")

    zip_path.unlink(missing_ok=True)
    return ttf_files


def download_from_google_fonts_repo(font_name, dest_dir):
    """Download font from google/fonts GitHub repo (FULL versions, not subsets)."""
    import urllib.request

    base = normalize_to_base(font_name)
    slug = base.lower().replace(' ', '')

    # Try to find the font in the google/fonts repo
    # Convention: ofl/<fontname>/ contains the TTF/variable files
    api_url = f"https://api.github.com/repos/google/fonts/contents/ofl/{slug}"

    try:
        req = urllib.request.Request(api_url, headers={
            'User-Agent': 'Masterminder/2.0',
            'Accept': 'application/vnd.github.v3+json'
        })
        with urllib.request.urlopen(req, timeout=15) as resp:
            contents = json.loads(resp.read().decode('utf-8'))
    except Exception:
        # Try apache license folder
        try:
            api_url = f"https://api.github.com/repos/google/fonts/contents/apache/{slug}"
            req = urllib.request.Request(api_url, headers={
                'User-Agent': 'Masterminder/2.0',
                'Accept': 'application/vnd.github.v3+json'
            })
            with urllib.request.urlopen(req, timeout=15) as resp:
                contents = json.loads(resp.read().decode('utf-8'))
        except Exception as e:
            print(f"    [!] Not found in Google Fonts repo: {e}")
            return []

    ttf_files = []
    for item in contents:
        name = item['name']
        # Prefer static TTFs, skip variable fonts
        if name.endswith('.ttf') and 'Variable' not in name and '[' not in name:
            download_url = item['download_url']
            out_path = dest_dir / name
            try:
                download_file(download_url, out_path, timeout=30)
                if out_path.stat().st_size > 50000:
                    ttf_files.append(out_path)
                else:
                    out_path.unlink()
            except Exception as e:
                print(f"    [!] Failed to download {name}: {e}")

    # If no static TTFs found, try the static/ subfolder
    if not ttf_files:
        try:
            static_url = f"{api_url.rstrip('/')}/static"
            req = urllib.request.Request(static_url, headers={
                'User-Agent': 'Masterminder/2.0',
                'Accept': 'application/vnd.github.v3+json'
            })
            with urllib.request.urlopen(req, timeout=15) as resp:
                static_contents = json.loads(resp.read().decode('utf-8'))

            for item in static_contents:
                name = item['name']
                if name.endswith('.ttf'):
                    out_path = dest_dir / name
                    try:
                        download_file(item['download_url'], out_path, timeout=30)
                        if out_path.stat().st_size > 50000:
                            ttf_files.append(out_path)
                    except Exception:
                        pass
        except Exception:
            pass

    return ttf_files


def check_cyrillic(ttf_path):
    """Check if a TTF file has Cyrillic glyphs. Returns (has_cyrillic, glyph_count)."""
    try:
        from fontTools.ttLib import TTFont
        font = TTFont(str(ttf_path))
        cmap = font.getBestCmap()
        found = sum(1 for cp in CYRILLIC_TEST if cp in cmap)
        return found >= 15, len(cmap)  # At least 15/21 Cyrillic chars
    except Exception:
        return False, 0


def install_fonts_windows(ttf_files):
    """Install TTF files for current user (no admin rights)."""
    import winreg
    import ctypes

    user_fonts = Path(os.path.expandvars(r'%LOCALAPPDATA%\Microsoft\Windows\Fonts'))
    user_fonts.mkdir(parents=True, exist_ok=True)

    installed = 0
    for ttf in ttf_files:
        ttf = Path(ttf)
        dest = user_fonts / ttf.name
        if not dest.exists() or dest.stat().st_size < ttf.stat().st_size:
            shutil.copy2(ttf, dest)
            with winreg.ConnectRegistry(None, winreg.HKEY_CURRENT_USER) as reg:
                with winreg.OpenKey(reg, r'Software\Microsoft\Windows NT\CurrentVersion\Fonts', 0, winreg.KEY_WRITE) as key:
                    font_label = ttf.stem.replace('-', ' ') + ' (TrueType)'
                    winreg.SetValueEx(key, font_label, 0, winreg.REG_SZ, str(dest))
            installed += 1

    if installed > 0:
        HWND_BROADCAST = 0xFFFF
        WM_FONTCHANGE = 0x001D
        ctypes.windll.user32.SendMessageW(HWND_BROADCAST, WM_FONTCHANGE, 0, 0)

    return installed


def replace_font_in_pptx(pptx_path, old_font, new_font):
    """Replace font name in PPTX XML."""
    temp_dir = Path(pptx_path).parent / '_pptx_font_fix'
    if temp_dir.exists():
        shutil.rmtree(temp_dir)
    temp_dir.mkdir()

    with zipfile.ZipFile(pptx_path, 'r') as z:
        z.extractall(temp_dir)

    changes = 0
    for root, dirs, files in os.walk(temp_dir):
        for fname in files:
            if fname.endswith(('.xml', '.rels')):
                fpath = os.path.join(root, fname)
                with open(fpath, 'r', encoding='utf-8') as f:
                    content = f.read()
                if old_font in content:
                    new_content = content.replace(old_font, new_font)
                    changes += content.count(old_font)
                    with open(fpath, 'w', encoding='utf-8') as f:
                        f.write(new_content)

    with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as z:
        for root, dirs, files in os.walk(temp_dir):
            for fname in files:
                fpath = os.path.join(root, fname)
                arcname = os.path.relpath(fpath, temp_dir)
                z.write(fpath, arcname)

    shutil.rmtree(temp_dir)
    return changes


def find_cyrillic_replacement(font_name):
    """Find best standard Windows replacement for a font without Cyrillic."""
    if font_name in NO_CYRILLIC_REPLACEMENTS:
        return NO_CYRILLIC_REPLACEMENTS[font_name]

    # Guess based on weight suffix
    base = normalize_to_base(font_name)
    if 'Light' in font_name:
        return 'Calibri Light'
    elif 'Bold' in font_name or 'SemiBold' in font_name or 'ExtraBold' in font_name:
        return 'Calibri'  # PowerPoint will apply bold style
    else:
        return 'Calibri'


def main():
    if len(sys.argv) < 2:
        print("Usage: python install-gamma-fonts.py <path-to-pptx>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    if platform.system() != 'Windows':
        print("[!] Font auto-install only supported on Windows.")
        sys.exit(0)

    print(f"=== Font Installer v2 ===")
    print(f"Analyzing: {pptx_path}\n")

    # Step 1: Extract fonts
    fonts = extract_fonts_from_pptx(pptx_path)
    print(f"Fonts in PPTX: {', '.join(f'{k} ({v}x)' for k, v in fonts.most_common())}\n")

    # Step 2: Find non-standard
    all_font_names = set(fonts.keys())
    non_standard = {f for f in all_font_names if not is_standard(f)}

    if not non_standard:
        print("All fonts are standard Windows fonts. Nothing to do.")
        return

    print(f"Non-standard fonts: {', '.join(sorted(non_standard))}\n")

    # Step 3: Download, verify Cyrillic, install
    temp_dir = Path(pptx_path).parent / '_font_temp'
    temp_dir.mkdir(exist_ok=True)

    fonts_replaced = {}  # old -> new (for fonts without Cyrillic)
    bases_processed = set()

    for font_name in sorted(non_standard):
        base = normalize_to_base(font_name)
        if base in bases_processed:
            continue
        bases_processed.add(base)

        print(f"--- {base} ---")

        # Download
        ttf_files = download_from_github(base, temp_dir)
        if not ttf_files:
            print(f"    Not on GitHub, trying Google Fonts repo...")
            ttf_files = download_from_google_fonts_repo(base, temp_dir)

        if not ttf_files:
            print(f"    [!] Could not download. Will replace with standard font.")
            for fn in all_font_names:
                if normalize_to_base(fn) == base:
                    replacement = find_cyrillic_replacement(fn)
                    fonts_replaced[fn] = replacement
            continue

        # Verify Cyrillic
        has_cyrillic = False
        for ttf in ttf_files:
            ok, glyph_count = check_cyrillic(ttf)
            if ok:
                has_cyrillic = True
                print(f"    {ttf.name}: {glyph_count} glyphs, Cyrillic: YES")
                break
            else:
                print(f"    {ttf.name}: {glyph_count} glyphs, Cyrillic: NO")

        if has_cyrillic:
            # Install full font
            count = install_fonts_windows(ttf_files)
            print(f"    Installed {count} files (full version with Cyrillic)")
        else:
            print(f"    [!] Font '{base}' has NO Cyrillic in any version.")
            for fn in all_font_names:
                if normalize_to_base(fn) == base:
                    replacement = find_cyrillic_replacement(fn)
                    fonts_replaced[fn] = replacement
                    print(f"    Will replace '{fn}' -> '{replacement}' in PPTX")

    # Step 4: Replace fonts without Cyrillic in PPTX XML
    if fonts_replaced:
        print(f"\n=== Replacing fonts without Cyrillic ===")
        for old, new in fonts_replaced.items():
            changes = replace_font_in_pptx(pptx_path, old, new)
            print(f"  '{old}' -> '{new}': {changes} replacements")

    # Cleanup
    shutil.rmtree(temp_dir, ignore_errors=True)

    print(f"\n=== Done ===")
    if fonts_replaced:
        print(f"Replaced {len(fonts_replaced)} font(s) without Cyrillic.")
    print("If PowerPoint is open — restart it to see the new fonts.")


if __name__ == '__main__':
    main()
